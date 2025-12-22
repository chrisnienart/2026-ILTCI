#!/usr/bin/env python3
"""
Script to apply markdown content to the ILTCI PowerPoint template.
This preserves all template styling, images, and backgrounds while adding your content.
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
import re
import yaml
import argparse
from pathlib import Path

def parse_markdown_slides(md_file):
    """Parse markdown file into individual slides."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    print(f"Original content length: {len(content)}")
    
    # Remove YAML frontmatter (everything between first --- and second ---)
    # Find the end of the frontmatter
    lines = content.split('\n')
    start_idx = 0
    end_idx = 0
    found_first = False
    
    for i, line in enumerate(lines):
        if line.strip() == '---':
            if not found_first:
                start_idx = i
                found_first = True
            else:
                end_idx = i
                break
    
    if end_idx > start_idx:
        # Remove frontmatter lines
        content_lines = lines[end_idx + 1:]
        content = '\n'.join(content_lines)
    
    print(f"Content after frontmatter removal length: {len(content)}")
    print("Content preview:", content[:200])
    
    # Split by slide separators (---) 
    # Use regex to split by --- on its own line
    slides = re.split(r'\n---\n', content)
    
    print(f"Number of slides found: {len(slides)}")
    
    parsed_slides = []
    for idx, slide in enumerate(slides):
        slide = slide.strip()
        if not slide:
            print(f"Slide {idx}: Empty, skipping")
            continue
        
        print(f"\n--- Parsing Slide {idx} ---")
        print(f"Slide content preview: {slide[:100]}")
        
        # Check for class directive
        is_title = '<!-- _class: title -->' in slide
        slide = slide.replace('<!-- _class: title -->', '').strip()
        
        print(f"Is title slide: {is_title}")
        
        # Parse slide content
        lines = slide.split('\n')
        title = ''
        subtitle_lines = []
        content_lines = []
        section_name = ''
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            print(f"  Processing line: {line[:50]}")
            
            if line.startswith('### ') and not section_name and is_title:
                # First h3 is section name on title slide
                section_name = line.replace('###', '').strip()
                print(f"  -> Section name: {section_name}")
            elif line.startswith('# '):
                title = line.replace('#', '').strip()
                print(f"  -> Title: {title}")
            elif line.startswith('## '):
                if is_title:
                    subtitle_lines.append(line.replace('##', '').strip())
                    print(f"  -> Subtitle: {line.replace('##', '').strip()}")
                else:
                    # For content slides, ## can also be a title
                    if not title:
                        title = line.replace('##', '').strip()
                        print(f"  -> Title (from ##): {title}")
            elif line.startswith('### '):
                if is_title:
                    subtitle_lines.append(line.replace('###', '').strip())
                    print(f"  -> Subtitle: {line.replace('###', '').strip()}")
                else:
                    content_lines.append(line)
                    print(f"  -> Content: {line}")
            else:
                content_lines.append(line)
                print(f"  -> Content: {line[:50]}")
        
        slide_data = {
            'is_title': is_title,
            'section_name': section_name,
            'title': title,
            'subtitle': '\n'.join(subtitle_lines),
            'content': '\n'.join(content_lines)
        }
        
        print(f"Parsed slide data: {slide_data}")
        parsed_slides.append(slide_data)
    
    print(f"\nTotal parsed slides: {len(parsed_slides)}")
    return parsed_slides

def remove_bullet(paragraph):
    """Remove bullet formatting from a paragraph."""
    pPr = paragraph._element.get_or_add_pPr()
    pPr.insert(0, OxmlElement('a:buNone'))

def add_numbering(paragraph, start_at=1):
    """Add automatic numbering to a paragraph."""
    pPr = paragraph._element.get_or_add_pPr()
    # Create buAutoNum element for numbering
    buAutoNum = OxmlElement('a:buAutoNum')
    buAutoNum.set('type', 'arabicPeriod')  # 1. 2. 3. style
    if start_at > 1:
        buAutoNum.set('startAt', str(start_at))
    pPr.insert(0, buAutoNum)

def add_formatted_text(paragraph, text):
    """
    Add text to a paragraph with markdown formatting support (bold, italic).
    Supports **bold**, *italic*, and ***bold+italic*** markdown syntax.
    """
    # Clear any existing text
    paragraph.text = ""
    
    # Pattern to match markdown formatting
    # Matches ***text*** (bold+italic), **text** (bold), or *text* (italic)
    pattern = r'(\*\*\*.*?\*\*\*|\*\*.*?\*\*|\*.*?\*)'
    
    # Split text by markdown patterns
    parts = re.split(pattern, text)
    
    for part in parts:
        if not part:
            continue
        
        # Check what type of formatting this part has
        if part.startswith('***') and part.endswith('***'):
            # Bold and italic
            run = paragraph.add_run()
            run.text = part[3:-3]
            run.font.bold = True
            run.font.italic = True
        elif part.startswith('**') and part.endswith('**'):
            # Bold
            run = paragraph.add_run()
            run.text = part[2:-2]
            run.font.bold = True
        elif part.startswith('*') and part.endswith('*'):
            # Italic
            run = paragraph.add_run()
            run.text = part[1:-1]
            run.font.italic = True
        else:
            # Regular text
            run = paragraph.add_run()
            run.text = part

def apply_to_template(template_file, md_file, output_file):
    """Apply markdown content to PowerPoint template."""
    print(f"\nLoading template: {template_file}")
    prs = Presentation(template_file)
    
    print(f"Template has {len(prs.slide_layouts)} layouts")
    print(f"Template has {len(prs.slides)} existing slides")
    
    parsed_slides = parse_markdown_slides(md_file)
    
    # Remove existing content slides (keep master/layouts)
    print("\nRemoving existing slides...")
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    print(f"Creating {len(parsed_slides)} new slides...")
    
    for idx, slide_data in enumerate(parsed_slides):
        print(f"\nCreating slide {idx + 1}...")
        print(f"  Title: {slide_data['title']}")
        print(f"  Is title slide: {slide_data['is_title']}")
        
        if slide_data['is_title']:
            # Use title slide layout (index 0)
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            print(f"  Using layout 0 (title slide), {len(slide.shapes)} shapes")
            
            # Debug: print all shapes
            for i, shape in enumerate(slide.shapes):
                print(f"    Shape {i}: {shape.name}, has_text_frame: {hasattr(shape, 'text_frame')}")
            
            # Try to set section name - look for the top text box
            if len(slide.shapes) > 2 and hasattr(slide.shapes[2], 'text_frame'):
                if slide_data['section_name']:
                    slide.shapes[2].text = slide_data['section_name']
                    print(f"  Set section name in shape 2")
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data['title']
                print(f"  Set title")
            
            # Set subtitle - try shape 1
            if len(slide.shapes) > 1 and hasattr(slide.shapes[1], 'text_frame'):
                slide.shapes[1].text = slide_data['subtitle']
                print(f"  Set subtitle in shape 1")
        else:
            # Use content slide layout (index 1)
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            print(f"  Using layout 1 (content slide), {len(slide.shapes)} shapes")
            
            # Debug: print all shapes
            for i, shape in enumerate(slide.shapes):
                print(f"    Shape {i}: {shape.name}, has_text_frame: {hasattr(shape, 'text_frame')}")
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data['title']
                print(f"  Set title")
            
            # Set content - find the content placeholder
            content_shape = None
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        content_shape = shape
                        break
            
            if content_shape and hasattr(content_shape, 'text_frame'):
                text_frame = content_shape.text_frame
                text_frame.clear()
                print(f"  Adding content to text frame...")
                
                # Parse content (bullets, numbered lists, etc.)
                for line in slide_data['content'].split('\n'):
                    line_stripped = line.strip()
                    if not line_stripped:
                        continue
                    
                    # Handle bullet points
                    if line_stripped.startswith('- '):
                        p = text_frame.add_paragraph()
                        add_formatted_text(p, line_stripped[2:])
                        p.level = 0
                        # Bullets are on by default in template, so no need to set
                        print(f"    Added bullet: {line_stripped[2:]}")
                    elif line_stripped.startswith('  - '):
                        p = text_frame.add_paragraph()
                        add_formatted_text(p, line_stripped[4:])
                        p.level = 1
                        # Bullets are on by default in template
                        print(f"    Added sub-bullet: {line_stripped[4:]}")
                    # Handle numbered lists (e.g., "1. ", "2. ")
                    elif re.match(r'^\d+\.\s+', line_stripped):
                        p = text_frame.add_paragraph()
                        # Extract the number and text
                        match = re.match(r'^(\d+)\.\s+(.*)$', line_stripped)
                        if match:
                            num = int(match.group(1))
                            text = match.group(2)
                        else:
                            num = 1
                            text = re.sub(r'^\d+\.\s+', '', line_stripped)
                        add_formatted_text(p, text)
                        p.level = 0
                        # Add automatic numbering
                        add_numbering(p, start_at=num)
                        print(f"    Added numbered item {num}: {text}")
                    else:
                        p = text_frame.add_paragraph()
                        add_formatted_text(p, line_stripped)
                        # Turn off bullets for regular text
                        remove_bullet(p)
                        print(f"    Added text: {line_stripped}")
    
    print(f"\nSaving presentation to {output_file}...")
    prs.save(output_file)
    print(f"✓ Presentation saved successfully!")
    print(f"  Total slides created: {len(prs.slides)}")

def load_config(config_path):
    """Load configuration from YAML file."""
    config_file = Path(config_path)
    if not config_file.exists():
        raise FileNotFoundError(f"Configuration file not found: {config_path}")
    
    with open(config_file, 'r', encoding='utf-8') as f:
        config = yaml.safe_load(f)
    
    return config

def parse_arguments():
    """Parse command-line arguments."""
    parser = argparse.ArgumentParser(
        description='Generate PowerPoint presentation from markdown content and template.'
    )
    
    # Config file path
    parser.add_argument(
        '--config',
        default='config.yaml',
        help='Path to configuration file (default: config.yaml)'
    )
    
    # Path overrides
    parser.add_argument(
        '--template',
        help='Path to PowerPoint template file (overrides config)'
    )
    
    parser.add_argument(
        '--content',
        help='Path to markdown content file (overrides config)'
    )
    
    parser.add_argument(
        '--output',
        help='Path to output PowerPoint file (overrides config)'
    )
    
    return parser.parse_args()

def main():
    """Main entry point for the script."""
    # Parse command-line arguments
    args = parse_arguments()
    
    # Load configuration
    try:
        config = load_config(args.config)
    except FileNotFoundError as e:
        print(f"Error: {e}")
        print(f"Please ensure the configuration file exists at: {args.config}")
        return 1
    except yaml.YAMLError as e:
        print(f"Error parsing configuration file: {e}")
        return 1
    
    # Extract paths from config with CLI overrides
    template = args.template or config['paths']['template']
    markdown = args.content or config['paths']['content']
    output = args.output or config['paths']['output']
    
    # Verify files exist
    template_path = Path(template)
    markdown_path = Path(markdown)
    
    if not template_path.exists():
        print(f"Error: Template file not found: {template}")
        return 1
    
    if not markdown_path.exists():
        print(f"Error: Content file not found: {markdown}")
        return 1
    
    # Create output directory if it doesn't exist
    output_path = Path(output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    print("=" * 60)
    print("ILTCI Presentation Generator")
    print("=" * 60)
    print(f"Configuration: {args.config}")
    print(f"Template:      {template}")
    print(f"Content:       {markdown}")
    print(f"Output:        {output}")
    print("=" * 60)
    
    try:
        apply_to_template(template, markdown, output)
    except Exception as e:
        print(f"\nError generating presentation: {e}")
        return 1
    
    print("\n" + "=" * 60)
    print("Done!")
    print("=" * 60)
    return 0

if __name__ == '__main__':
    exit(main())
